#!/usr/bin/env python3
"""
Convert all 6 workflow HTML slides to a single PowerPoint presentation.
Uses Playwright to screenshot each slide at 1440x810, then assembles with python-pptx.
"""

import subprocess
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SLIDES_DIR = Path(__file__).parent
OUTPUT = SLIDES_DIR / "Nova_Workflow_Maps.pptx"
SCREENSHOTS_DIR = SLIDES_DIR / "_screenshots"

SLIDE_FILES = [
    ("1-recruiter-workflow.html", "Recruiter Workflow"),
    ("2-marketing-manager.html", "Marketing Manager Workflow"),
    ("3-designer-workflow.html", "Designer Workflow"),
    ("4-agency-workflow.html", "Agency Portal Workflow"),
    ("5-system-pipeline.html", "System Pipeline — All 6 Stages"),
    ("6-team-process.html", "Team Process — End-to-End Flow"),
]

# 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def screenshot_slides():
    """Use Playwright to screenshot each HTML file."""
    SCREENSHOTS_DIR.mkdir(exist_ok=True)

    for filename, _ in SLIDE_FILES:
        html_path = SLIDES_DIR / filename
        png_path = SCREENSHOTS_DIR / filename.replace(".html", ".png")

        if not html_path.exists():
            print(f"  SKIP {filename} — not found")
            continue

        print(f"  Capturing {filename}...")
        # Use playwright screenshot CLI
        file_url = f"file://{html_path.resolve()}"
        subprocess.run([
            "npx", "playwright", "screenshot",
            "--browser", "chromium",
            "--viewport-size", "1440,810",
            "--full-page",
            file_url,
            str(png_path),
        ], check=True, capture_output=True)
        print(f"    → {png_path.name}")


def build_pptx():
    """Assemble screenshots into a PPTX."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Blank layout
    blank_layout = prs.slide_layouts[6]

    for filename, title in SLIDE_FILES:
        png_path = SCREENSHOTS_DIR / filename.replace(".html", ".png")
        if not png_path.exists():
            print(f"  SKIP {title} — no screenshot")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add full-bleed image
        slide.shapes.add_picture(
            str(png_path),
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )

        # Add slide notes with title
        slide.notes_slide.notes_text_frame.text = title

        print(f"  Added: {title}")

    prs.save(str(OUTPUT))
    print(f"\n✅ Saved: {OUTPUT}")
    print(f"   {len(prs.slides)} slides")


if __name__ == "__main__":
    print("Step 1: Screenshot HTML slides...")
    screenshot_slides()
    print("\nStep 2: Build PowerPoint...")
    build_pptx()
